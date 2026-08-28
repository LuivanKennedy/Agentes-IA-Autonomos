"""
Módulo de Ferramentas - Projeto Agentes IA Autônomos
Responsável pela geração de documentos, leitura local, buscas na web e APIs.

Destaques desta versão:
- Excel em MODO DINÂMICO: projeção 100% por fórmulas nativas (FV) ligadas à aba 'Premissas'.
- 3 cenários (pessimista / base / otimista) + benchmark CDI no mesmo gráfico.
- Formatação condicional na carteira (ícones de tendência, barras de dados, escala de cores).
- Dupla compatibilidade: as fórmulas são gravadas junto com seus VALORES em cache,
  permitindo leitura por pandas/openpyxl sem precisar abrir o Excel.
"""

import os
import re
import time
import json
import shutil
import zipfile
import requests
import pandas as pd
import urllib.parse
import xml.etree.ElementTree as ET
from fpdf import FPDF
from bs4 import BeautifulSoup

# --- Imports para Word ---
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH

# --- Imports para PowerPoint ---
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor

# --- Imports para Excel ---
from openpyxl import Workbook, load_workbook
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.chart.label import DataLabelList, DataLabel
from openpyxl.formatting.rule import IconSetRule, ColorScaleRule, DataBarRule
from openpyxl.utils import get_column_letter
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.workbook.properties import CalcProperties

# ==========================================
# 0. SISTEMA DE CACHE E PESQUISA (VALIDADE DE 24H)
# ==========================================
_cache_busca = {}
_cache_pagina = {}


def carregar_cache():
    """Carrega os caches do disco. Se tiverem mais de 24h, são apagados."""
    global _cache_busca, _cache_pagina
    agora = time.time()

    def ler_arquivo_cache(nome_arquivo):
        if os.path.exists(nome_arquivo):
            idade_arquivo = agora - os.path.getmtime(nome_arquivo)  # 86400 = 24h
            if idade_arquivo > 86400:
                print(f"🧹 Limpando cache antigo (mais de 24h): {nome_arquivo}")
                os.remove(nome_arquivo)
                return {}
            try:
                with open(nome_arquivo, "r", encoding="utf-8") as f:
                    return json.load(f)
            except Exception:
                return {}
        return {}

    _cache_busca = ler_arquivo_cache("cache_busca.json")
    _cache_pagina = ler_arquivo_cache("cache_pagina.json")


def salvar_cache():
    """Salva os caches no disco para persistência entre execuções."""
    with open("cache_busca.json", "w", encoding="utf-8") as f:
        json.dump(_cache_busca, f, ensure_ascii=False, indent=2)
    with open("cache_pagina.json", "w", encoding="utf-8") as f:
        json.dump(_cache_pagina, f, ensure_ascii=False, indent=2)


carregar_cache()


# ==========================================
# 1. FERRAMENTAS DE CRIAÇÃO (DOCUMENTOS)
# ==========================================

def criar_apresentacao_com_ia(dados_apresentacao: dict, nome_arquivo: str = "apresentacao.pptx") -> str:
    print(f"🖥️ Montando Apresentação Avançada: {nome_arquivo}")
    prs = Presentation()
    slides_dados = dados_apresentacao.get("slides", [])

    for i, slide_info in enumerate(slides_dados):
        titulo = slide_info.get("titulo", f"Slide {i+1}")
        texto = slide_info.get("texto", "")
        prompt_imagem = slide_info.get("prompt_imagem", "")

        print(f"   - Processando slide {i+1}: {titulo}")

        if i == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = titulo
            title.text_frame.paragraphs[0].font.size = PptxPt(44)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(0, 51, 102)
            subtitle.text = texto
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title_shape = slide.shapes.title
            title_shape.text = titulo
            title_shape.text_frame.paragraphs[0].font.size = PptxPt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = PptxRGBColor(0, 51, 102)

            caminho_imagem = None
            if prompt_imagem and prompt_imagem.lower() != "nenhuma":
                for _ in range(3):
                    try:
                        prompt_cod = urllib.parse.quote(prompt_imagem)
                        url_img = (f"https://image.pollinations.ai/prompt/{prompt_cod}"
                                   f"?width=800&height=600&nologo=true&model=flux")
                        response = requests.get(url_img, timeout=45)
                        if response.status_code == 200:
                            caminho_imagem = f"temp_slide_{i}.jpg"
                            with open(caminho_imagem, "wb") as f:
                                f.write(response.content)
                            break
                    except Exception:
                        time.sleep(2)

            if caminho_imagem:
                slide.shapes.add_picture(caminho_imagem, PptxInches(0.5), PptxInches(2.0),
                                         width=PptxInches(4.5))
                txBox = slide.shapes.add_textbox(PptxInches(5.2), PptxInches(2.0),
                                                 PptxInches(4.3), PptxInches(4.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = texto
                p.font.size = PptxPt(18)
                try:
                    os.remove(caminho_imagem)
                except Exception:
                    pass
            else:
                txBox = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(2.0),
                                                 PptxInches(8.0), PptxInches(4.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = texto
                p.font.size = PptxPt(20)

    caminho_final = os.path.join(os.getcwd(), nome_arquivo)
    prs.save(caminho_final)
    return f"Apresentação '{nome_arquivo}' criada com sucesso!"


def _adicionar_texto_com_negrito(paragrafo, texto):
    partes = re.split(r'(\*\*.*?\*\*)', texto)
    for parte in partes:
        if parte.startswith('**') and parte.endswith('**') and len(parte) > 4:
            run = paragrafo.add_run(parte[2:-2])
            run.bold = True
        else:
            paragrafo.add_run(parte)


def criar_documento_word(texto: str, nome_arquivo: str = "relatorio.docx",
                         imagens_para_inserir: list = None) -> str:
    print(f"📝 Gerando documento Word profissional: {nome_arquivo}")

    texto = _normalizar_quebras(texto)
    texto = re.sub(r'\s*[—–]\s*', ' - ', texto)
    texto = re.sub(r'\s*-\s*,', ',', texto)
    texto = re.sub(r'\s*-\s*\.', '.', texto)

    doc = Document()
    linhas = texto.splitlines()

    imagens_colocadas = 0
    imagens_inseridas_por_tag = []

    for linha in linhas:
        linha_limpa = linha.strip()
        if not linha_limpa:
            continue
        if linha_limpa.replace('-', '') == '' or linha_limpa.replace('*', '') == '':
            continue

        # == CAÇADOR DE IMAGENS ==
        match_tag = re.search(r'\[IMAGEM:\s*(.+?)\]', linha_limpa, re.IGNORECASE)
        match_md = re.search(r'!\[.*?\]\((.+?)\)', linha_limpa)

        img_path = None
        if match_tag:
            img_path = match_tag.group(1).strip()
        elif match_md:
            img_path = match_md.group(1).strip()

        if img_path:
            if os.path.exists(img_path):
                print(f"   - Inserindo imagem no meio do texto: {img_path}")
                p_img = doc.add_paragraph()
                p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p_img.paragraph_format.space_before = DocxPt(12)
                p_img.paragraph_format.space_after = DocxPt(12)
                p_img.add_run().add_picture(img_path, width=Inches(5.5))
                imagens_colocadas += 1
                imagens_inseridas_por_tag.append(img_path)
            else:
                print(f"   - Aviso: Imagem '{img_path}' solicitada não foi encontrada.")
            continue

        # == TÍTULOS E SUBTÍTULOS ==
        match_heading = re.match(r'^(#+)\s+(.*)', linha_limpa)
        if match_heading:
            nivel_word = min(len(match_heading.group(1)), 9)
            texto_titulo = match_heading.group(2).replace('**', '')
            h = doc.add_heading(texto_titulo, level=nivel_word)
            h.paragraph_format.space_before = DocxPt(18)
            h.paragraph_format.space_after = DocxPt(12)
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER if nivel_word == 1 else WD_ALIGN_PARAGRAPH.JUSTIFY
            continue

        # == LISTAS (Bolinhas) ==
        if re.match(r'^[-*•]\s+', linha_limpa):
            p = doc.add_paragraph(style='List Bullet')
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            _adicionar_texto_com_negrito(p, re.sub(r'^[-*•]\s+', '', linha_limpa))
            continue

        # == LISTAS (Números) ==
        if re.match(r'^\d+\.\s+', linha_limpa):
            p = doc.add_paragraph(style='List Number')
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            _adicionar_texto_com_negrito(p, re.sub(r'^\d+\.\s+', '', linha_limpa))
            continue

        # == PARÁGRAFOS NORMAIS ==
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        p.paragraph_format.first_line_indent = Cm(1.25)
        _adicionar_texto_com_negrito(p, linha_limpa)

        for palavra in ['importante:', 'atenção:', 'conclusão:']:
            if palavra in linha_limpa.lower():
                if len(p.runs) > 0:
                    p.runs[0].font.color.rgb = DocxRGBColor(200, 0, 0)
                break

    # == PLANO B (FALLBACK DE IMAGENS) ==
    if imagens_para_inserir:
        imagens_faltantes = [img for img in imagens_para_inserir if img not in imagens_inseridas_por_tag]
        if imagens_faltantes:
            doc.add_paragraph()
            h_fallback = doc.add_heading('Anexos e Ilustrações Adicionais', level=1)
            h_fallback.alignment = WD_ALIGN_PARAGRAPH.CENTER

            for img in imagens_faltantes:
                if os.path.exists(img):
                    print(f"   - Inserindo imagem via Fallback (Final do Doc): {img}")
                    p_fb = doc.add_paragraph()
                    p_fb.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p_fb.paragraph_format.space_before = DocxPt(12)
                    p_fb.add_run().add_picture(img, width=Inches(5.5))
                    p_leg = doc.add_paragraph(f"Ilustração: {img}")
                    p_leg.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    imagens_colocadas += 1

    caminho_final = os.path.join(os.getcwd(), nome_arquivo)
    doc.save(caminho_final)
    return f"Documento Word '{nome_arquivo}' salvo com sucesso! {imagens_colocadas} imagens inseridas."


class RelatorioPDF(FPDF):
    def header(self):
        self.set_font('Arial', 'B', 12)
        self.set_text_color(0, 51, 102)
        self.cell(0, 10, 'Relatório Gerado por IA Autônoma', 0, 1, 'R')
        self.line(10, 20, 200, 20)
        self.ln(10)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f'Página {self.page_no()}', 0, 0, 'C')


def criar_pdf(texto: str, nome_arquivo: str = "documento.pdf") -> str:
    print(f"📄 Gerando PDF executivo: {nome_arquivo}")
    pdf = RelatorioPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    texto_limpo = _normalizar_quebras(texto).replace('**', '').replace('##', '').replace('#', '')
    for linha in texto_limpo.split('\n'):
        linha = linha.strip()
        if not linha:
            pdf.ln(5)
            continue
        if len(linha) < 60 and not linha.endswith('.'):
            pdf.set_font('Arial', 'B', 14)
            pdf.set_text_color(0, 51, 102)
            pdf.cell(0, 10, linha.encode('latin-1', 'replace').decode('latin-1'), ln=True)
            pdf.ln(2)
        else:
            pdf.set_font('Arial', '', 11)
            pdf.set_text_color(0, 0, 0)
            pdf.multi_cell(0, 6, linha.encode('latin-1', 'replace').decode('latin-1'))

    caminho_final = os.path.join(os.getcwd(), nome_arquivo)
    pdf.output(caminho_final)
    return f"PDF '{nome_arquivo}' criado com sucesso."


# ==========================================
# 1.1 HELPERS DO EXCEL
# ==========================================

PALAVRAS_MOEDA = ['valor', 'preço', 'preco', 'custo', 'total', 'salário', 'salario', 'r$',
                  'aporte', 'projeção', 'projecao', 'retorno', 'patrimônio', 'patrimonio',
                  'cdi', 'saldo', 'montante', 'investimento', 'cotação', 'cotacao',
                  'pessimista', 'otimista', 'base']
PALAVRAS_PCT = ['porcentagem', 'percentual', '%', 'taxa', 'yield', 'rentabilidade',
                'dy', 'dividend', 'peso', 'alocação', 'alocacao']
PALAVRAS_TEMPO = ['ano', 'anos', 'mês', 'mes', 'data', 'período', 'periodo', 'year']
PALAVRAS_ID = ['ativo', 'ticker', 'papel', 'nome', 'categoria', 'classe', 'setor',
               'descrição', 'descricao', 'item', 'grupo', 'segmento']

AZUL = '1F4E78'
ZEBRA = 'F2F7FB'
AMARELO_EDITAVEL = 'FFF2CC'


def _normalizar_quebras(texto: str) -> str:
    """CORREÇÃO CRÍTICA: a IA costuma enviar a sequência literal '\\n' (barra + letra n)
    em vez de uma quebra de linha real, o que fazia todo o insight virar uma única
    linha gigante e ser cortado na célula."""
    if not isinstance(texto, str):
        return ""
    texto = texto.replace('\\r\\n', '\n').replace('\\n', '\n').replace('\\t', ' ')
    texto = texto.replace('\r\n', '\n').replace('\r', '\n')
    return texto


def _para_float(valor, padrao=None):
    """Converte com segurança valores vindos da IA ('R$ 1.000,00', '13,5%', 0.135)."""
    if isinstance(valor, (int, float)):
        return float(valor)
    if not isinstance(valor, str):
        return padrao
    v = valor.replace("R$", "").replace("%", "").strip()
    if "," in v and "." in v:
        v = v.replace(".", "").replace(",", ".") if v.rfind(",") > v.rfind(".") else v.replace(",", "")
    elif "," in v:
        v = v.replace(",", ".")
    try:
        return float(v)
    except ValueError:
        return padrao


def _normalizar_taxa(valor, padrao):
    """Aceita 13.5 (percentual) ou 0.135 (decimal) e devolve sempre decimal."""
    v = _para_float(valor, None)
    if v is None:
        return padrao
    return v / 100 if v > 1 else v


def _detectar_col_tempo(df):
    for c in df.columns:
        if str(c).lower().strip() in PALAVRAS_TEMPO:
            return c
    return None


def _classificar_coluna(nome_col):
    n = str(nome_col).lower()
    if n.strip() in PALAVRAS_TEMPO:
        return "tempo"
    if any(p in n for p in PALAVRAS_PCT):
        return "pct"
    if any(p in n for p in PALAVRAS_MOEDA):
        return "moeda"
    if any(p in n for p in PALAVRAS_ID):
        return "id"
    return "outro"


def _ordenar_colunas(df):
    """A API do Gemini não garante a ordem das chaves do JSON, o que embaralhava as
    colunas. Reordena logicamente: tempo -> identificadores -> moeda -> pct -> resto.
    Também remove colunas com conteúdo duplicado."""
    tempo, ids, moeda, pct, outros = [], [], [], [], []
    for c in df.columns:
        tipo = _classificar_coluna(c)
        {"tempo": tempo, "id": ids, "moeda": moeda, "pct": pct}.get(tipo, outros).append(c)

    ids.sort(key=lambda c: 0 if str(c).lower().strip() in ('ativo', 'ticker', 'papel') else 1)

    df = df[tempo + ids + moeda + pct + outros]

    colunas_manter, vistos = [], []
    for c in df.columns:
        assinatura = tuple(df[c].astype(str).tolist())
        if assinatura in vistos:
            print(f"   - Removendo coluna duplicada: {c}")
            continue
        vistos.append(assinatura)
        colunas_manter.append(c)

    return df[colunas_manter]


def _formatar_tabela(ws, df, primeira_linha=1):
    """Aplica estilos, converte números e ajusta larguras.
    Retorna (colunas_moeda, colunas_porcentagem, col_tempo)."""
    cor_fundo = PatternFill(start_color=AZUL, end_color=AZUL, fill_type='solid')
    fonte_branca = Font(bold=True, color='FFFFFF')
    zebra = PatternFill(start_color=ZEBRA, end_color=ZEBRA, fill_type='solid')
    borda_fina = Border(
        left=Side(style='thin', color='D9D9D9'),
        right=Side(style='thin', color='D9D9D9'),
        top=Side(style='thin', color='D9D9D9'),
        bottom=Side(style='thin', color='D9D9D9')
    )

    col_tempo = _detectar_col_tempo(df)

    colunas_moeda, colunas_porcentagem = [], []
    for idx, col_name in enumerate(df.columns, start=1):
        tipo = _classificar_coluna(col_name)
        if tipo == "pct":
            colunas_porcentagem.append(idx)
        elif tipo == "moeda":
            colunas_moeda.append(idx)

    for cell in ws[primeira_linha]:
        cell.fill = cor_fundo
        cell.font = fonte_branca
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = borda_fina
    ws.row_dimensions[primeira_linha].height = 30
    ws.freeze_panes = f"A{primeira_linha + 1}"
    ws.auto_filter.ref = (f"A{primeira_linha}:"
                          f"{get_column_letter(len(df.columns))}{primeira_linha + len(df)}")

    for col_idx in range(1, len(df.columns) + 1):
        tamanho_max = len(str(ws.cell(row=primeira_linha, column=col_idx).value or "")) or 10

        for row_idx in range(primeira_linha + 1, primeira_linha + len(df) + 1):
            cell = ws.cell(row=row_idx, column=col_idx)
            cell.border = borda_fina
            if row_idx % 2 == 0:
                cell.fill = zebra

            # --- SISTEMA ANTI-ALUCINAÇÃO DA IA (padrão US/BR) ---
            if isinstance(cell.value, str) and (col_idx in colunas_moeda or col_idx in colunas_porcentagem):
                novo_valor = _para_float(cell.value, None)
                if novo_valor is not None:
                    if col_idx in colunas_porcentagem and novo_valor > 1:
                        novo_valor = novo_valor / 100
                    cell.value = novo_valor

            if col_idx in colunas_moeda and isinstance(cell.value, (int, float)):
                cell.number_format = '"R$" #,##0.00'
                cell.alignment = Alignment(horizontal='right')
            elif col_idx in colunas_porcentagem and isinstance(cell.value, (int, float)):
                cell.number_format = '0.00%'
                cell.alignment = Alignment(horizontal='center')
            elif col_tempo is not None and col_idx == df.columns.get_loc(col_tempo) + 1:
                cell.alignment = Alignment(horizontal='center')

            if cell.value is not None:
                if col_idx in colunas_moeda and isinstance(cell.value, (int, float)):
                    tam_texto = len(f"R$ {cell.value:,.2f}")
                else:
                    tam_texto = len(str(cell.value))
                tamanho_max = max(tamanho_max, tam_texto)

        ws.column_dimensions[get_column_letter(col_idx)].width = min(max(tamanho_max + 3, 12), 42)

    return colunas_moeda, colunas_porcentagem, col_tempo


def _rotular_apenas_ultimo_ponto(serie, indice_ultimo):
    """Mostra o valor SOMENTE no último ponto da linha (valor final da projeção),
    em vez de poluir o gráfico com 30 rótulos sobrepostos."""
    dl = DataLabel(idx=indice_ultimo)
    dl.showVal = True
    dl.showSerName = False
    dl.showCatName = False
    dl.showLegendKey = False
    dl.showBubbleSize = False
    dl.showPercent = False

    serie.dLbls = DataLabelList()
    serie.dLbls.dLbl.append(dl)
    serie.dLbls.showVal = False
    serie.dLbls.showSerName = False
    serie.dLbls.showCatName = False
    serie.dLbls.showLegendKey = False
    serie.dLbls.showBubbleSize = False
    serie.dLbls.showPercent = False


def _escrever_bloco_texto(ws, linha_inicial, titulo, texto, largura_colunas=10):
    """Escreve texto longo em MÚLTIPLAS linhas (uma por parágrafo), com mesclagem
    horizontal e altura calculada. Sem corte e sem linhas vazias inúteis."""
    col_fim = get_column_letter(max(largura_colunas, 8))

    cel_titulo = ws.cell(row=linha_inicial, column=1, value=titulo)
    cel_titulo.font = Font(bold=True, size=13, color="FFFFFF")
    cel_titulo.fill = PatternFill(start_color=AZUL, end_color=AZUL, fill_type='solid')
    cel_titulo.alignment = Alignment(horizontal='left', vertical='center', indent=1)
    ws.merge_cells(f"A{linha_inicial}:{col_fim}{linha_inicial}")
    ws.row_dimensions[linha_inicial].height = 24

    largura_util = 0
    for i in range(1, max(largura_colunas, 8) + 1):
        dim = ws.column_dimensions[get_column_letter(i)].width
        largura_util += dim if dim else 10
    largura_util = max(int(largura_util), 60)

    paragrafos = [p.strip() for p in texto.split('\n') if p.strip()]

    linha = linha_inicial + 1
    for p in paragrafos:
        cel = ws.cell(row=linha, column=1, value=p)
        negrito = bool(re.match(r'^\d+\.\s', p)) or (p.isupper() and len(p) > 3)
        cel.font = Font(bold=negrito, size=11, color=AZUL if negrito else "000000")
        cel.alignment = Alignment(wrap_text=True, vertical="center", horizontal="left", indent=1)
        ws.merge_cells(f"A{linha}:{col_fim}{linha}")

        linhas_visuais = max(1, -(-len(p) // largura_util))
        ws.row_dimensions[linha].height = 16 * linhas_visuais + 4
        linha += 1

    return linha


# ==========================================
# 1.2 MOTOR DE CÁLCULO ESPELHO (CACHE DE FÓRMULAS)
# ==========================================

def _fv(taxa_anual: float, anos: float, aporte_mensal: float, aporte_inicial: float) -> float:
    """Replica EXATAMENTE a função FV do Excel usada na planilha:
    =-FV((1+taxa)^(1/12)-1 ; anos*12 ; aporte_mensal ; aporte_inicial)
    Capitalização mensal equivalente à taxa anual informada."""
    i = (1 + taxa_anual) ** (1 / 12) - 1
    n = anos * 12
    if i == 0:
        return aporte_inicial + aporte_mensal * n
    fator = (1 + i) ** n
    return aporte_inicial * fator + aporte_mensal * (fator - 1) / i


def _calcular_cache_premissas(p: dict) -> dict:
    """Valores em cache das células calculadas da aba 'Premissas'."""
    anos = p["anos"]
    taxa_pes = p["taxa_base"] - p["desvio_pes"]
    taxa_oti = p["taxa_base"] + p["desvio_oti"]

    total_aportado = p["aporte_inicial"] + p["aporte_mensal"] * anos * 12
    base = _fv(p["taxa_base"], anos, p["aporte_mensal"], p["aporte_inicial"])
    cdi = _fv(p["taxa_cdi"], anos, p["aporte_mensal"], p["aporte_inicial"])

    return {
        "B13": taxa_pes,
        "B14": taxa_oti,
        "B18": total_aportado,
        "B19": base,
        "B20": base - total_aportado,
        "B21": base / ((1 + p["inflacao"]) ** anos),
        "B22": base - cdi,
    }


def _calcular_cache_projecao(p: dict) -> dict:
    """Valores em cache de toda a tabela de projeção (colunas B a G)."""
    cache = {}
    taxas = {
        "B": p["taxa_base"] - p["desvio_pes"],
        "C": p["taxa_base"],
        "D": p["taxa_base"] + p["desvio_oti"],
        "E": p["taxa_cdi"],
    }
    for ano in range(1, p["anos"] + 1):
        linha = ano + 1
        for col, taxa in taxas.items():
            cache[f"{col}{linha}"] = _fv(taxa, ano, p["aporte_mensal"], p["aporte_inicial"])
        cache[f"F{linha}"] = p["aporte_inicial"] + p["aporte_mensal"] * ano * 12
        cache[f"G{linha}"] = cache[f"C{linha}"] / ((1 + p["inflacao"]) ** ano)
    return cache


# ------------------------------------------
# 1.3 INJEÇÃO DO CACHE NO XML DO .XLSX
# ------------------------------------------

NS_MAIN = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
NS_REL_DOC = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _mapear_abas_xlsx(zf: zipfile.ZipFile) -> dict:
    """Descobre qual arquivo XML interno corresponde a cada nome de aba."""
    wb_xml = ET.fromstring(zf.read("xl/workbook.xml"))
    rels_xml = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))

    rels = {}
    for rel in rels_xml:
        alvo = rel.get("Target", "")
        alvo = alvo[1:] if alvo.startswith("/") else (
            alvo if alvo.startswith("xl/") else f"xl/{alvo}")
        rels[rel.get("Id")] = alvo.replace("xl/xl/", "xl/")

    mapa = {}
    for sheet in wb_xml.find(f"{{{NS_MAIN}}}sheets"):
        rid = sheet.get(f"{{{NS_REL_DOC}}}id")
        if rid in rels:
            mapa[sheet.get("name")] = rels[rid]
    return mapa


def _patch_sheet_xml(conteudo: bytes, valores: dict) -> bytes:
    """Insere <v>valor</v> nas células que possuem <f>formula</f>."""
    ET.register_namespace("", NS_MAIN)
    raiz = ET.fromstring(conteudo)
    tag_f, tag_v = f"{{{NS_MAIN}}}f", f"{{{NS_MAIN}}}v"

    for celula in raiz.iter(f"{{{NS_MAIN}}}c"):
        ref = celula.get("r")
        if ref not in valores or celula.find(tag_f) is None:
            continue

        valor = valores[ref]
        if valor is None or (isinstance(valor, float) and valor != valor):  # NaN
            continue

        for antigo in celula.findall(tag_v):
            celula.remove(antigo)

        elem_v = ET.SubElement(celula, tag_v)
        elem_v.text = repr(round(float(valor), 10)) if isinstance(valor, float) else str(valor)

        # célula de fórmula numérica não deve carregar t="s"/t="str"
        if celula.get("t") in ("s", "str"):
            del celula.attrib["t"]

    return ET.tostring(raiz, encoding="UTF-8", xml_declaration=True)


def _injetar_cache_formulas(caminho: str, cache_por_aba: dict) -> bool:
    """Reescreve o .xlsx adicionando os resultados calculados junto às fórmulas.
    Assim a planilha funciona nos dois mundos:
      - Excel/LibreOffice: recalcula ao vivo pelas fórmulas;
      - pandas/openpyxl: lê os números do cache, sem precisar abrir o Excel."""
    if not cache_por_aba:
        return False

    temporario = caminho + ".tmp"
    try:
        with zipfile.ZipFile(caminho, "r") as origem:
            mapa = _mapear_abas_xlsx(origem)
            destinos = {mapa[a]: v for a, v in cache_por_aba.items() if a in mapa}

            with zipfile.ZipFile(temporario, "w", zipfile.ZIP_DEFLATED) as saida:
                for item in origem.infolist():
                    dados = origem.read(item.filename)
                    if item.filename in destinos:
                        dados = _patch_sheet_xml(dados, destinos[item.filename])
                    saida.writestr(item, dados)

        shutil.move(temporario, caminho)
        total = sum(len(v) for v in cache_por_aba.values())
        print(f"   - Cache de fórmulas injetado: {total} células calculadas gravadas.")
        return True
    except Exception as e:
        print(f"      ⚠️ Falha ao injetar cache de fórmulas: {e}")
        if os.path.exists(temporario):
            os.remove(temporario)
        return False


# ------------------------------------------
# 1.4 ABA DE PREMISSAS EDITÁVEIS
# ------------------------------------------

# Endereços fixos usados pelas fórmulas da projeção
CEL = {
    "aporte_inicial": "Premissas!$B$4",
    "aporte_mensal": "Premissas!$B$5",
    "anos": "Premissas!$B$6",
    "taxa_base": "Premissas!$B$9",
    "taxa_cdi": "Premissas!$B$10",
    "desvio_pes": "Premissas!$B$11",
    "desvio_oti": "Premissas!$B$12",
    "taxa_pes": "Premissas!$B$13",
    "taxa_oti": "Premissas!$B$14",
    "inflacao": "Premissas!$B$15",
}


def _criar_aba_premissas(wb, params):
    """Cria a aba 'Premissas' com células AMARELAS editáveis. Alterar qualquer valor
    aqui recalcula toda a projeção e os 3 cenários automaticamente no Excel."""
    ws = wb.create_sheet("Premissas")
    ws.sheet_view.showGridLines = False
    ws.column_dimensions['A'].width = 42
    ws.column_dimensions['B'].width = 20
    ws.column_dimensions['C'].width = 60

    fill_titulo = PatternFill(start_color=AZUL, end_color=AZUL, fill_type='solid')
    fill_edit = PatternFill(start_color=AMARELO_EDITAVEL, end_color=AMARELO_EDITAVEL, fill_type='solid')
    borda = Border(*(Side(style='thin', color='BFBFBF'),) * 4)

    def secao(linha, texto):
        for col in range(1, 4):
            c = ws.cell(row=linha, column=col)
            c.fill = fill_titulo
            c.font = Font(bold=True, color="FFFFFF", size=12)
        ws.cell(row=linha, column=1, value=texto).alignment = Alignment(indent=1, vertical='center')
        ws.row_dimensions[linha].height = 22

    def campo(linha, rotulo, valor, formato, ajuda, editavel=True, formula=None):
        ws.cell(row=linha, column=1, value=rotulo).font = Font(bold=True, size=11)
        c = ws.cell(row=linha, column=2, value=formula if formula else valor)
        c.number_format = formato
        c.alignment = Alignment(horizontal='center')
        c.border = borda
        if editavel:
            c.fill = fill_edit
            c.font = Font(bold=True, size=11, color="7F6000")
        else:
            c.font = Font(bold=True, size=11, color=AZUL, italic=True)
        obs = ws.cell(row=linha, column=3, value=ajuda)
        obs.font = Font(size=9, color="808080", italic=True)
        obs.alignment = Alignment(indent=1)

    t = ws.cell(row=1, column=1, value="⚙️ PAINEL DE PREMISSAS EDITÁVEIS")
    t.font = Font(bold=True, size=16, color=AZUL)
    ws.cell(row=2, column=1,
            value="Altere APENAS as células amarelas. Toda a projeção e os gráficos recalculam sozinhos."
            ).font = Font(size=10, italic=True, color="C00000")

    secao(3, "APORTES E PRAZO")
    campo(4, "Aporte Inicial (R$)", params["aporte_inicial"], '"R$" #,##0.00',
          "Capital investido no momento zero.")
    campo(5, "Aporte Mensal (R$)", params["aporte_mensal"], '"R$" #,##0.00',
          "Valor investido todo mês, com dividendos reinvestidos.")
    campo(6, "Prazo (anos)", params["anos"], '0',
          "Horizonte da projeção. A tabela é gerada com o prazo máximo definido na criação.")

    secao(8, "TAXAS DE RETORNO (a.a.)")
    campo(9, "Cenário Base — Carteira", params["taxa_base"], '0.00%',
          "Retorno nominal esperado da carteira recomendada.")
    campo(10, "Benchmark — CDI", params["taxa_cdi"], '0.00%',
          "Taxa livre de risco usada como comparação.")
    campo(11, "Desvio Cenário Pessimista", params["desvio_pes"], '0.00%',
          "Quanto subtrair da taxa base (ex.: 4% => base menos 4 p.p.).")
    campo(12, "Desvio Cenário Otimista", params["desvio_oti"], '0.00%',
          "Quanto somar à taxa base.")
    campo(13, "→ Taxa Pessimista (calculada)", None, '0.00%',
          "Fórmula automática: Base - Desvio.", editavel=False, formula="=B9-B11")
    campo(14, "→ Taxa Otimista (calculada)", None, '0.00%',
          "Fórmula automática: Base + Desvio.", editavel=False, formula="=B9+B12")
    campo(15, "Inflação estimada (IPCA)", params["inflacao"], '0.00%',
          "Usada para calcular o patrimônio em valores de hoje (poder de compra).")

    secao(17, "RESULTADOS NO FIM DO PRAZO")
    ws.cell(row=18, column=1, value="Total Aportado (nominal)").font = Font(bold=True)
    c18 = ws.cell(row=18, column=2, value="=B4+(B5*B6*12)")
    c18.number_format = '"R$" #,##0.00'
    c18.font = Font(bold=True, color=AZUL)
    ws.cell(row=18, column=3, value="Soma do capital do próprio bolso, sem rendimento."
            ).font = Font(size=9, italic=True, color="808080")

    ws.cell(row=19, column=1, value="Patrimônio Final — Cenário Base").font = Font(bold=True)
    c19 = ws.cell(row=19, column=2, value="=-FV((1+B9)^(1/12)-1,B6*12,B5,B4)")
    c19.number_format = '"R$" #,##0.00'
    c19.font = Font(bold=True, size=12, color="006100")
    ws.cell(row=19, column=3, value="Juros compostos com capitalização mensal."
            ).font = Font(size=9, italic=True, color="808080")

    ws.cell(row=20, column=1, value="Juros Ganhos (Base)").font = Font(bold=True)
    c20 = ws.cell(row=20, column=2, value="=B19-B18")
    c20.number_format = '"R$" #,##0.00'
    c20.font = Font(bold=True, color="006100")

    ws.cell(row=21, column=1, value="Poder de Compra Hoje (Base)").font = Font(bold=True)
    c21 = ws.cell(row=21, column=2, value="=B19/((1+B15)^B6)")
    c21.number_format = '"R$" #,##0.00'
    c21.font = Font(bold=True, color="9C6500")
    ws.cell(row=21, column=3, value="Patrimônio final descontada a inflação do período."
            ).font = Font(size=9, italic=True, color="808080")

    ws.cell(row=22, column=1, value="Vantagem sobre o CDI").font = Font(bold=True)
    c22 = ws.cell(row=22, column=2, value="=B19-(-FV((1+B10)^(1/12)-1,B6*12,B5,B4))")
    c22.number_format = '"R$" #,##0.00'
    c22.font = Font(bold=True, color="C00000")

    # Validações: impede taxas absurdas / negativas por digitação errada
    dv_taxa = DataValidation(type="decimal", operator="between", formula1="-0.5", formula2="1",
                             allow_blank=False, showErrorMessage=True,
                             error="Informe a taxa em formato decimal entre -50% e 100%.",
                             errorTitle="Taxa inválida")
    ws.add_data_validation(dv_taxa)
    for ref in ("B9", "B10", "B11", "B12", "B15"):
        dv_taxa.add(ws[ref])

    dv_prazo = DataValidation(type="whole", operator="between", formula1="1", formula2="60",
                              showErrorMessage=True, error="O prazo deve ficar entre 1 e 60 anos.",
                              errorTitle="Prazo inválido")
    ws.add_data_validation(dv_prazo)
    dv_prazo.add(ws["B6"])

    return ws


def _montar_projecao_com_formulas(ws, anos, titulo_colunas=None):
    """Escreve a tabela de projeção usando FÓRMULAS FV nativas do Excel, ligadas à
    aba Premissas. Gera 3 cenários + CDI, tudo recalculável pelo usuário."""
    cabecalhos = titulo_colunas or ["Ano", "Cenário Pessimista", "Cenário Base",
                                    "Cenário Otimista", "Benchmark CDI",
                                    "Total Aportado", "Base (valor de hoje)"]
    ws.append(cabecalhos)

    taxas = {
        2: CEL["taxa_pes"],
        3: CEL["taxa_base"],
        4: CEL["taxa_oti"],
        5: CEL["taxa_cdi"],
    }

    for i in range(1, anos + 1):
        linha = i + 1
        ws.cell(row=linha, column=1, value=i)
        for col, taxa in taxas.items():
            ws.cell(
                row=linha, column=col,
                value=(f"=-FV((1+{taxa})^(1/12)-1,A{linha}*12,"
                       f"{CEL['aporte_mensal']},{CEL['aporte_inicial']})")
            )
        ws.cell(row=linha, column=6,
                value=f"={CEL['aporte_inicial']}+({CEL['aporte_mensal']}*A{linha}*12)")
        ws.cell(row=linha, column=7,
                value=f"=C{linha}/((1+{CEL['inflacao']})^A{linha})")

    return cabecalhos


def _aplicar_formatacao_condicional_carteira(ws, df_at, ultima_linha):
    """Ícones de tendência no Dividend Yield, barras de dados no Aporte e
    escala de cores no Peso da Carteira."""
    mapa = {}
    for i, c in enumerate(df_at.columns, start=1):
        nome = str(c).lower()
        letra = get_column_letter(i)
        faixa = f"{letra}2:{letra}{ultima_linha}"
        if 'yield' in nome or nome.strip() in ('dy',) or 'dividend' in nome:
            mapa['dy'] = faixa
        elif 'aporte' in nome:
            mapa['aporte'] = faixa
        elif 'peso' in nome or 'alocação' in nome or 'alocacao' in nome:
            mapa['peso'] = faixa
        elif 'preço' in nome or 'preco' in nome:
            mapa['preco'] = faixa

    if 'dy' in mapa:
        # 3 setas: verde (DY alto), amarela (médio), vermelha (baixo/sem provento)
        regra = IconSetRule('3Arrows', 'num', ['0.02', '0.06'], showValue=True)
        ws.conditional_formatting.add(mapa['dy'], regra)

    if 'aporte' in mapa:
        ws.conditional_formatting.add(
            mapa['aporte'],
            DataBarRule(start_type='num', start_value=0, end_type='max',
                        color="1F4E78", showValue=True)
        )

    if 'peso' in mapa:
        ws.conditional_formatting.add(
            mapa['peso'],
            ColorScaleRule(start_type='min', start_color='FFFFFF',
                           end_type='max', end_color='9BC2E6')
        )

    return mapa


def criar_planilha_excel(dados: list = None, insights: str = "", nome_arquivo: str = "dados.xlsx",
                         ativos: list = None, titulo_grafico: str = None,
                         premissas: dict = None) -> str:
    """
    dados          -> tabela estática (usada só quando NÃO há 'premissas').
    insights       -> texto analítico da IA (múltiplas linhas, sem corte).
    ativos         -> (opcional) carteira recomendada; gera a aba 'Carteira Recomendada'.
    titulo_grafico -> (opcional) título do gráfico principal.
    premissas      -> (RECOMENDADO para projeções financeiras) dicionário com:
                      aporte_inicial, aporte_mensal, anos, taxa_base, taxa_cdi,
                      desvio_pessimista, desvio_otimista, inflacao.
                      Quando enviado, a planilha é gerada com FÓRMULAS EXCEL VIVAS
                      e 3 cenários, dispensando a lista 'dados'.

    Abas: 'Premissas', 'Projeção e Análise', 'Carteira Recomendada' e 'Dashboard'.
    """
    print(f"📊 Gerando Excel Inteligente: {nome_arquivo}")

    modo_dinamico = bool(premissas)
    if not modo_dinamico and not dados:
        return "Erro: envie 'premissas' (modo dinâmico) ou 'dados' (modo estático)."

    wb = Workbook()
    ws = wb.active
    ws.title = "Projeção e Análise"

    df = None
    colunas_moeda, colunas_porcentagem, col_tempo, col_agrupador = [], [], None, None
    params_dinamicos = None
    cache_carteira = {}

    # =========================================================
    # MODO A: PROJEÇÃO DINÂMICA COM FÓRMULAS E 3 CENÁRIOS
    # =========================================================
    if modo_dinamico:
        params = {
            "aporte_inicial": _para_float(premissas.get("aporte_inicial"), 0.0) or 0.0,
            "aporte_mensal": _para_float(premissas.get("aporte_mensal"), 0.0) or 0.0,
            "anos": int(_para_float(premissas.get("anos"), 30) or 30),
            "taxa_base": _normalizar_taxa(premissas.get("taxa_base"), 0.12),
            "taxa_cdi": _normalizar_taxa(premissas.get("taxa_cdi"), 0.10),
            "desvio_pes": _normalizar_taxa(premissas.get("desvio_pessimista"), 0.04),
            "desvio_oti": _normalizar_taxa(premissas.get("desvio_otimista"), 0.03),
            "inflacao": _normalizar_taxa(premissas.get("inflacao"), 0.045),
        }
        params["anos"] = max(1, min(params["anos"], 60))
        params_dinamicos = params

        # Força o Excel a recalcular ao abrir, mesmo com os valores já em cache
        wb.calculation = CalcProperties(calcId=124519, fullCalcOnLoad=True)

        _criar_aba_premissas(wb, params)
        cabecalhos = _montar_projecao_com_formulas(ws, params["anos"])

        df = pd.DataFrame(columns=cabecalhos, index=range(params["anos"]))
        colunas_moeda = [2, 3, 4, 5, 6, 7]
        col_tempo = "Ano"
        ultima_linha = params["anos"] + 1

        borda = Border(*(Side(style='thin', color='D9D9D9'),) * 4)
        for cell in ws[1]:
            cell.fill = PatternFill(start_color=AZUL, end_color=AZUL, fill_type='solid')
            cell.font = Font(bold=True, color='FFFFFF')
            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            cell.border = borda
        ws.row_dimensions[1].height = 32
        ws.freeze_panes = "B2"

        for row in range(2, ultima_linha + 1):
            for col in range(1, 8):
                c = ws.cell(row=row, column=col)
                c.border = borda
                if row % 2 == 0:
                    c.fill = PatternFill(start_color=ZEBRA, end_color=ZEBRA, fill_type='solid')
                if col == 1:
                    c.alignment = Alignment(horizontal='center')
                    c.font = Font(bold=True)
                else:
                    c.number_format = '"R$" #,##0.00'
                    c.alignment = Alignment(horizontal='right')

        larguras = [8, 20, 20, 20, 20, 20, 22]
        for i, w in enumerate(larguras, start=1):
            ws.column_dimensions[get_column_letter(i)].width = w

        # Destaque visual: o cenário base é a coluna principal
        for row in range(2, ultima_linha + 1):
            ws.cell(row=row, column=3).font = Font(bold=True, color="006100")

        aviso = ws.cell(row=ultima_linha + 1, column=1,
                        value="🔗 Valores calculados por fórmula. Edite a aba 'Premissas' "
                              "para recalcular tudo automaticamente.")
        aviso.font = Font(italic=True, size=10, color="C00000")
        ws.merge_cells(start_row=ultima_linha + 1, start_column=1,
                       end_row=ultima_linha + 1, end_column=7)

    # =========================================================
    # MODO B: TABELA ESTÁTICA (compatibilidade retroativa)
    # =========================================================
    else:
        df = pd.DataFrame(dados)

        if list(df.columns) == [0, 1]:
            df.columns = ["Categoria", "Valor"]
        elif list(df.columns) == [0, 1, 2]:
            df.columns = ["Grupo", "Categoria", "Valor"]

        df = _ordenar_colunas(df)

        for c in df.columns:
            if str(c).lower() in ['contexto', 'grupo', 'tipo', 'classificação']:
                col_agrupador = c
                break

        col_tempo_pre = _detectar_col_tempo(df)
        if col_tempo_pre is not None:
            df["__ordem__"] = pd.to_numeric(df[col_tempo_pre], errors='coerce')
            df = df.sort_values(by="__ordem__", kind='mergesort', na_position='last')
            df = df.drop(columns=["__ordem__"]).reset_index(drop=True)
        elif col_agrupador and df[col_agrupador].nunique() > 1:
            df = df.sort_values(by=col_agrupador, kind='mergesort').reset_index(drop=True)

        for r in dataframe_to_rows(df, index=False, header=True):
            ws.append(r)

        colunas_moeda, colunas_porcentagem, col_tempo = _formatar_tabela(ws, df)
        ultima_linha = len(df) + 1

    # --- INSIGHTS EM MÚLTIPLAS LINHAS (SEM CORTE) ---
    insights_txt = _normalizar_quebras(insights or "")
    if insights_txt:
        insights_txt = re.sub(r'(R\$)\s*\n\s*([\d.,]+)', r'\1 \2', insights_txt)
        insights_txt = re.sub(r'(?m)^[ \t]*(\d+\.)[ \t]+', r'\n\1 ', insights_txt)
        insights_txt = re.sub(r'\n{3,}', '\n\n', insights_txt).strip()
    else:
        insights_txt = "Nenhum insight gerado."

    linha_insights = ultima_linha + (3 if modo_dinamico else 2)
    _escrever_bloco_texto(ws, linha_insights, "💡 INSIGHTS E ANÁLISE DA IA",
                          insights_txt, largura_colunas=max(len(df.columns), 8))

    # --- ABA: CARTEIRA RECOMENDADA (COM FORMATAÇÃO CONDICIONAL) ---
    ws2 = None
    df_at = None
    col_aporte = None
    if ativos:
        try:
            df_at = _ordenar_colunas(pd.DataFrame(ativos))
            ws2 = wb.create_sheet("Carteira Recomendada")
            for r in dataframe_to_rows(df_at, index=False, header=True):
                ws2.append(r)

            _formatar_tabela(ws2, df_at)
            ultima_linha_at = len(df_at) + 1

            _aplicar_formatacao_condicional_carteira(ws2, df_at, ultima_linha_at)

            for i, c in enumerate(df_at.columns, start=1):
                if 'aporte' in str(c).lower():
                    col_aporte = i
                    break

            col_rotulo_total = 1
            for i, c in enumerate(df_at.columns, start=1):
                if _classificar_coluna(c) == "id":
                    col_rotulo_total = i
                    break

            if col_aporte:
                linha_total = ultima_linha_at + 1
                estilo_total = PatternFill(start_color=AZUL, end_color=AZUL, fill_type='solid')

                for i in range(1, len(df_at.columns) + 1):
                    cel = ws2.cell(row=linha_total, column=i)
                    cel.fill = estilo_total
                    cel.font = Font(bold=True, color="FFFFFF")

                ws2.cell(row=linha_total, column=col_rotulo_total, value="TOTAL DOS APORTES")

                letra = get_column_letter(col_aporte)
                cel_total = ws2.cell(row=linha_total, column=col_aporte,
                                     value=f"=SUM({letra}2:{letra}{ultima_linha_at})")
                cel_total.number_format = '"R$" #,##0.00'
                cel_total.alignment = Alignment(horizontal='right')

                # Cache do total, para leitura sem abrir o Excel
                soma_aportes = sum(
                    v for v in (_para_float(x, None) for x in df_at.iloc[:, col_aporte - 1])
                    if v is not None
                )
                cache_carteira = {f"{letra}{linha_total}": soma_aportes}

                # Legenda dos ícones de tendência
                linha_legenda = linha_total + 2
                leg = ws2.cell(row=linha_legenda, column=1,
                               value="Legenda dos ícones no Dividend Yield:  ▲ verde = DY acima de 6% "
                                     "(forte gerador de renda)   ▬ amarelo = entre 2% e 6%   "
                                     "▼ vermelho = abaixo de 2% (foco em ganho de capital)")
                leg.font = Font(size=9, italic=True, color="595959")
                ws2.merge_cells(start_row=linha_legenda, start_column=1,
                                end_row=linha_legenda, end_column=max(len(df_at.columns), 4))
        except Exception as e:
            print(f"      ⚠️ Falha ao criar a aba de ativos: {e}")

    # --- ABA: DASHBOARD (GRÁFICOS ISOLADOS) ---
    try:
        wsd = wb.create_sheet("Dashboard")
        wsd.sheet_view.showGridLines = False
        cel_t = wsd.cell(row=1, column=1, value="📈 PAINEL DE GRÁFICOS")
        cel_t.font = Font(bold=True, size=16, color=AZUL)
        wsd.cell(row=2, column=1,
                 value="Os gráficos estão vinculados às fórmulas: mudar a aba 'Premissas' "
                       "atualiza as curvas na hora."
                 ).font = Font(size=10, italic=True, color="808080")

        ancora = 4

        if col_tempo is not None and colunas_moeda:
            chart = LineChart()
            chart.style = 12
            chart.title = titulo_grafico or "Projeção por Cenário x CDI"
            chart.y_axis.title = "Patrimônio Acumulado (R$)"
            chart.x_axis.title = str(col_tempo)
            chart.y_axis.numFmt = '"R$" #,##0'
            chart.height = 15
            chart.width = 32

            for c_idx in colunas_moeda:
                chart.add_data(Reference(ws, min_col=c_idx, min_row=1, max_row=ultima_linha),
                               titles_from_data=True)

            col_tempo_idx = (1 if modo_dinamico else df.columns.get_loc(col_tempo) + 1)
            chart.set_categories(Reference(ws, min_col=col_tempo_idx, min_row=2, max_row=ultima_linha))

            passo = max(1, (ultima_linha - 1) // 15)
            chart.x_axis.tickLblSkip = passo
            chart.x_axis.tickMarkSkip = passo
            chart.legend.position = 'b'

            # Paleta: pessimista vermelho, base verde grosso, otimista azul, CDI cinza
            paleta = ["C00000", "006100", "2E75B6", "808080", "BFBFBF", "9C6500"]
            for i, s in enumerate(chart.series):
                s.smooth = False
                cor = paleta[i % len(paleta)]
                s.graphicalProperties.line.solidFill = cor
                s.graphicalProperties.line.width = 38000 if i == 1 else 22000
                if i in (0, 2):
                    s.graphicalProperties.line.dashStyle = "dash"
                _rotular_apenas_ultimo_ponto(s, ultima_linha - 2)

            wsd.add_chart(chart, f"A{ancora}")
            ancora += 32

            # Gráfico 2: aportes x juros compostos (só no modo dinâmico)
            if modo_dinamico:
                comp = BarChart()
                comp.type = "col"
                comp.grouping = "stacked"
                comp.overlap = 100
                comp.style = 11
                comp.title = "Capital Aportado x Patrimônio Projetado (Cenário Base)"
                comp.y_axis.numFmt = '"R$" #,##0'
                comp.height = 12
                comp.width = 32
                comp.add_data(Reference(ws, min_col=6, min_row=1, max_row=ultima_linha),
                              titles_from_data=True)
                comp.add_data(Reference(ws, min_col=3, min_row=1, max_row=ultima_linha),
                              titles_from_data=True)
                comp.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ultima_linha))
                comp.legend.position = 'b'
                wsd.add_chart(comp, f"A{ancora}")
                ancora += 26

        elif col_agrupador:
            grupos = df[col_agrupador].unique()
            col_agrupador_idx = df.columns.get_loc(col_agrupador) + 1
            col_rotulo = 1
            for i in range(1, len(df.columns) + 1):
                if i != col_agrupador_idx and i not in colunas_moeda and i not in colunas_porcentagem:
                    col_rotulo = i
                    break
            col_grafico = (colunas_moeda[0] if colunas_moeda
                           else (colunas_porcentagem[0] if colunas_porcentagem else 2))

            coluna_letra = ['A', 'K', 'U']
            for idx_grp, grupo in enumerate(grupos):
                indices = df.index[df[col_agrupador] == grupo].tolist()
                chart = BarChart()
                chart.type = "col"
                chart.style = 10 + (idx_grp % 10)
                chart.title = str(grupo)
                chart.height = 11
                chart.width = 14
                chart.legend = None
                chart.dataLabels = DataLabelList()
                chart.dataLabels.showVal = True
                chart.add_data(Reference(ws, min_col=col_grafico, min_row=indices[0] + 2,
                                         max_row=indices[-1] + 2), titles_from_data=False)
                chart.set_categories(Reference(ws, min_col=col_rotulo, min_row=indices[0] + 2,
                                               max_row=indices[-1] + 2))
                wsd.add_chart(chart, f"{coluna_letra[idx_grp % 3]}{ancora + (idx_grp // 3) * 24}")
            ancora += ((len(grupos) - 1) // 3 + 1) * 24

        else:
            chart = BarChart()
            chart.type = "col"
            chart.style = 10
            chart.title = titulo_grafico or "Análise Visual de Dados"
            chart.height = 12
            chart.width = 24
            chart.dataLabels = DataLabelList()
            chart.dataLabels.showVal = True
            col_grafico = colunas_moeda[0] if colunas_moeda else 2
            chart.add_data(Reference(ws, min_col=col_grafico, min_row=1,
                                     max_row=ultima_linha, max_col=col_grafico),
                           titles_from_data=True)
            chart.set_categories(Reference(ws, min_col=1, min_row=2, max_row=ultima_linha))
            wsd.add_chart(chart, f"A{ancora}")
            ancora += 26

        if ws2 is not None and df_at is not None and col_aporte:
            col_nome_ativo = 1
            for i, c in enumerate(df_at.columns, start=1):
                if str(c).lower().strip() in ('ativo', 'ticker', 'papel'):
                    col_nome_ativo = i
                    break

            bar = BarChart()
            bar.type = "bar"
            bar.style = 11
            bar.title = "Distribuição do Aporte Mensal por Ativo"
            bar.x_axis.title = "Aporte (R$)"
            bar.height = 13
            bar.width = 24
            bar.legend = None
            bar.add_data(Reference(ws2, min_col=col_aporte, min_row=1, max_row=len(df_at) + 1),
                         titles_from_data=True)
            bar.set_categories(Reference(ws2, min_col=col_nome_ativo, min_row=2, max_row=len(df_at) + 1))
            bar.dataLabels = DataLabelList()
            bar.dataLabels.showVal = True
            wsd.add_chart(bar, f"A{ancora}")
    except Exception as e:
        print(f"      ⚠️ Falha ao montar o Dashboard: {e}")

    # Reordena as abas: Premissas primeiro (é o painel de controle)
    if modo_dinamico:
        try:
            ordem = ["Premissas", "Projeção e Análise", "Carteira Recomendada", "Dashboard"]
            wb._sheets.sort(key=lambda s: ordem.index(s.title) if s.title in ordem else 99)
            wb.active = 0
        except Exception:
            pass

    caminho_final = os.path.join(os.getcwd(), nome_arquivo)
    wb.save(caminho_final)

    # --- GRAVA OS RESULTADOS JUNTO ÀS FÓRMULAS (dupla compatibilidade) ---
    cache_por_aba = {}
    if params_dinamicos:
        cache_por_aba["Premissas"] = _calcular_cache_premissas(params_dinamicos)
        cache_por_aba["Projeção e Análise"] = _calcular_cache_projecao(params_dinamicos)
    if cache_carteira:
        cache_por_aba["Carteira Recomendada"] = cache_carteira

    cache_ok = _injetar_cache_formulas(caminho_final, cache_por_aba)

    if modo_dinamico:
        resumo = _calcular_cache_premissas(params_dinamicos)
        return (f"Planilha '{nome_arquivo}' gerada em MODO DINÂMICO: aba 'Premissas' editável, "
                f"projeção por fórmulas FV, 3 cenários e formatação condicional na carteira. "
                f"Patrimônio final no cenário base: R$ {resumo['B19']:,.2f} "
                f"(total aportado R$ {resumo['B18']:,.2f}; poder de compra de hoje "
                f"R$ {resumo['B21']:,.2f}). "
                f"{'Valores gravados em cache, legíveis sem abrir o Excel.' if cache_ok else ''}")

    return (f"Planilha '{nome_arquivo}' gerada com sucesso (modo estático). "
            f"Colunas reordenadas e gráficos na aba 'Dashboard'.")


def gerar_imagem(prompt: str, nome_arquivo: str = "imagem_gerada.jpg") -> str:
    print(f"🎨 Gerando imagem: {nome_arquivo}")
    try:
        prompt_turbinado = (f"{prompt}, award-winning photography, highly detailed, "
                            f"8k resolution, ultra-realistic, cinematic lighting")
        prompt_codificado = urllib.parse.quote(prompt_turbinado)
        url_img = (f"https://image.pollinations.ai/prompt/{prompt_codificado}"
                   f"?width=1024&height=576&nologo=true&model=flux")
        response = requests.get(url_img, timeout=45)

        if response.status_code == 200:
            caminho_final = os.path.join(os.getcwd(), nome_arquivo)
            with open(caminho_final, "wb") as f:
                f.write(response.content)
            return f"Imagem gerada e salva com sucesso em: {caminho_final}"
        return "Erro: Falha no servidor de imagens."
    except Exception as e:
        return f"Erro ao gerar imagem: {e}"


# ==========================================
# 2. FERRAMENTAS DE PESQUISA E LEITURA
# ==========================================

def _ler_excel_robusto(caminho: str) -> str:
    """Lê planilhas priorizando VALORES calculados. Se encontrar fórmulas sem cache,
    avisa explicitamente em vez de devolver '=-FV(...)' como se fosse dado."""
    partes = []
    wb_val = load_workbook(caminho, data_only=True)
    wb_frm = load_workbook(caminho, data_only=False)

    for nome in wb_val.sheetnames:
        ws_v, ws_f = wb_val[nome], wb_frm[nome]
        linhas, sem_cache = [], 0

        for lin_v, lin_f in zip(ws_v.iter_rows(), ws_f.iter_rows()):
            celulas = []
            for cv, cf in zip(lin_v, lin_f):
                valor = cv.value
                formula = isinstance(cf.value, str) and cf.value.startswith("=")
                if valor is None and formula:
                    sem_cache += 1
                    celulas.append("[fórmula sem valor em cache]")
                elif isinstance(valor, float):
                    celulas.append(f"{valor:,.2f}")
                else:
                    celulas.append("" if valor is None else str(valor))
            if any(c.strip() for c in celulas):
                linhas.append(" | ".join(celulas).rstrip(" |"))

        cabecalho = f"--- ABA: {nome} ---"
        if sem_cache:
            cabecalho += (f"\n[AVISO: {sem_cache} célula(s) de fórmula sem valor gravado. "
                          f"Abra e salve o arquivo no Excel para materializá-las.]")
        partes.append(cabecalho + "\n" + ("\n".join(linhas) if linhas else "(aba vazia)"))

    return "Conteúdo da Planilha Excel:\n" + "\n\n".join(partes)


def ler_arquivo(caminho: str) -> str:
    print(f"📖 Lendo arquivo: {caminho}")
    if not os.path.exists(caminho):
        return f"Erro: O arquivo '{caminho}' não foi encontrado no sistema."

    ext = caminho.split('.')[-1].lower()
    try:
        if ext in ('txt', 'csv'):
            with open(caminho, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext in ('xlsx', 'xls'):
            if ext == 'xlsx':
                try:
                    return _ler_excel_robusto(caminho)
                except Exception as e:
                    print(f"      ⚠️ Leitura robusta falhou ({e}); usando pandas.")
            planilhas = pd.read_excel(caminho, sheet_name=None)
            partes = [f"--- ABA: {nome} ---\n{d.to_string()}" for nome, d in planilhas.items()]
            return "Conteúdo da Planilha Excel:\n" + "\n\n".join(partes)

        elif ext == 'docx':
            doc = Document(caminho)
            return "\n".join([p.text for p in doc.paragraphs])

        elif ext == 'pdf':
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(caminho)
                texto = "\n".join([p.extract_text() for p in reader.pages if p.extract_text()])
                return texto if texto else "O PDF não contém texto legível."
            except ImportError:
                return "Erro interno: instale a biblioteca 'PyPDF2' (pip install PyPDF2)."
        else:
            return f"Formato de arquivo '.{ext}' não suportado para leitura direta."
    except Exception as e:
        return f"Erro ao tentar ler o arquivo: {e}"


def buscar_web(query: str) -> str:
    print(f"🔍 Buscando na Web: {query}")
    if query in _cache_busca:
        return _cache_busca[query]

    try:
        from ddgs import DDGS
        resultados = DDGS().text(query, max_results=5)
        if not resultados:
            return "Nenhum resultado encontrado."

        texto_resultado = "\n".join(
            [f"Título: {r['title']}\nLink: {r['href']}\nResumo: {r['body']}\n" for r in resultados]
        )
        _cache_busca[query] = texto_resultado
        salvar_cache()
        return texto_resultado
    except ImportError:
        return "Erro: Instale a biblioteca ddgs (pip install ddgs)."
    except Exception as e:
        return f"Erro na busca: {e}"


def abrir_pagina(url: str) -> str:
    print(f"🌐 Lendo página: {url}")
    if url in _cache_pagina:
        return _cache_pagina[url]

    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
        resposta = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(resposta.text, 'html.parser')

        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.extract()

        texto = soup.get_text(separator=' ', strip=True)[:8000]
        _cache_pagina[url] = texto
        salvar_cache()
        return texto
    except Exception as e:
        return f"Erro ao acessar a página: {e}"


def cotacao_moeda(moeda: str) -> str:
    moeda = moeda.upper().strip()
    print(f"💱 Buscando cotação: {moeda}")
    try:
        url = f"https://economia.awesomeapi.com.br/last/{moeda}-BRL"
        dados = requests.get(url, timeout=10).json()
        chave = f"{moeda}BRL"
        return (f"Cotação {moeda}/BRL: R$ {dados[chave]['bid']} "
                f"(Atualizado em: {dados[chave]['create_date']})")
    except Exception as e:
        return f"Erro ao buscar cotação. Verifique se a moeda é válida (USD, EUR, BTC). Erro: {e}"


# ==========================================
# 3. REGISTRO GERAL NO DICIONÁRIO
# ==========================================
FERRAMENTAS = {
    "criar_apresentacao_com_ia": criar_apresentacao_com_ia,
    "criar_documento_word": criar_documento_word,
    "criar_pdf": criar_pdf,
    "criar_planilha_excel": criar_planilha_excel,
    "gerar_imagem": gerar_imagem,

    "ler_arquivo": ler_arquivo,
    "buscar_web": buscar_web,
    "abrir_pagina": abrir_pagina,
    "cotacao_moeda": cotacao_moeda
}
